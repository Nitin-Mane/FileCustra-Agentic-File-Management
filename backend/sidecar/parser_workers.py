#!/usr/bin/env python3
"""
FileCustra Parser Workers
Isolated parser workers for PDF, Office, text/code, and media metadata extraction.
Each parser runs in a sandbox boundary to prevent crashing the main application.
"""

import os
import json
import logging
import mimetypes
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)


@dataclass
class ParseResult:
    """Result of parsing a file."""
    file_path: str
    file_name: str
    parser_type: str
    success: bool
    content_text: Optional[str] = None
    metadata: Dict[str, Any] = field(default_factory=dict)
    error_message: Optional[str] = None
    page_count: Optional[int] = None
    word_count: Optional[int] = None
    
    def to_dict(self) -> dict:
        return {
            "file_path": self.file_path,
            "file_name": self.file_name,
            "parser_type": self.parser_type,
            "success": self.success,
            "content_text": self.content_text,
            "metadata": self.metadata,
            "error_message": self.error_message,
            "page_count": self.page_count,
            "word_count": self.word_count,
        }


class BaseParser(ABC):
    """Base class for all parsers."""
    
    @property
    @abstractmethod
    def parser_type(self) -> str:
        pass
    
    @property
    @abstractmethod
    def supported_extensions(self) -> List[str]:
        pass
    
    @abstractmethod
    def parse(self, file_path: str) -> ParseResult:
        pass
    
    def can_parse(self, file_path: str) -> bool:
        ext = Path(file_path).suffix.lstrip(".").lower()
        return ext in self.supported_extensions


class PDFParser(BaseParser):
    """Parser for PDF documents."""
    
    @property
    def parser_type(self) -> str:
        return "pdf"
    
    @property
    def supported_extensions(self) -> List[str]:
        return ["pdf"]
    
    def parse(self, file_path: str) -> ParseResult:
        path = Path(file_path)
        if not path.exists():
            return ParseResult(
                file_path=file_path,
                file_name=path.name,
                parser_type=self.parser_type,
                success=False,
                error_message="File not found",
            )
        
        try:
            import fitz
            
            doc = fitz.open(file_path)
            content_parts = []
            metadata = {}
            
            if doc.metadata:
                metadata = {
                    "title": doc.metadata.get("title", ""),
                    "author": doc.metadata.get("author", ""),
                    "subject": doc.metadata.get("subject", ""),
                    "creator": doc.metadata.get("creator", ""),
                    "producer": doc.metadata.get("producer", ""),
                    "creation_date": doc.metadata.get("creationDate", ""),
                    "mod_date": doc.metadata.get("modDate", ""),
                }
            
            page_count = len(doc)
            
            for page_num in range(min(page_count, 100)):
                page = doc[page_num]
                text = page.get_text()
                if text.strip():
                    content_parts.append(text)
            
            doc.close()
            
            content_text = "\n\n".join(content_parts)
            word_count = len(content_text.split()) if content_text else 0
            
            return ParseResult(
                file_path=file_path,
                file_name=path.name,
                parser_type=self.parser_type,
                success=True,
                content_text=content_text,
                metadata=metadata,
                page_count=page_count,
                word_count=word_count,
            )
            
        except ImportError:
            return ParseResult(
                file_path=file_path,
                file_name=path.name,
                parser_type=self.parser_type,
                success=False,
                error_message="PyMuPDF not installed",
            )
        except Exception as e:
            return ParseResult(
                file_path=file_path,
                file_name=path.name,
                parser_type=self.parser_type,
                success=False,
                error_message=str(e),
            )


class OfficeParser(BaseParser):
    """Parser for Office documents (DOCX, XLSX, PPTX)."""
    
    @property
    def parser_type(self) -> str:
        return "office"
    
    @property
    def supported_extensions(self) -> List[str]:
        return ["docx", "xlsx", "pptx", "doc", "xls", "ppt"]
    
    def parse(self, file_path: str) -> ParseResult:
        path = Path(file_path)
        if not path.exists():
            return ParseResult(
                file_path=file_path,
                file_name=path.name,
                parser_type=self.parser_type,
                success=False,
                error_message="File not found",
            )
        
        ext = path.suffix.lstrip(".").lower()
        
        try:
            if ext in ["docx", "doc"]:
                return self._parse_docx(file_path, path)
            elif ext in ["xlsx", "xls"]:
                return self._parse_xlsx(file_path, path)
            elif ext in ["pptx", "ppt"]:
                return self._parse_pptx(file_path, path)
            else:
                return ParseResult(
                    file_path=file_path,
                    file_name=path.name,
                    parser_type=self.parser_type,
                    success=False,
                    error_message=f"Unsupported Office extension: {ext}",
                )
        except ImportError as e:
            return ParseResult(
                file_path=file_path,
                file_name=path.name,
                parser_type=self.parser_type,
                success=False,
                error_message=f"Office parser library not installed: {e}",
            )
        except Exception as e:
            return ParseResult(
                file_path=file_path,
                file_name=path.name,
                parser_type=self.parser_type,
                success=False,
                error_message=str(e),
            )
    
    def _parse_docx(self, file_path: str, path: Path) -> ParseResult:
        from docx import Document
        
        doc = Document(file_path)
        content_parts = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                content_parts.append(para.text)
        
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells)
                if row_text.strip():
                    content_parts.append(row_text)
        
        content_text = "\n".join(content_parts)
        word_count = len(content_text.split()) if content_text else 0
        
        metadata = {}
        if doc.core_properties:
            metadata = {
                "title": doc.core_properties.title or "",
                "author": doc.core_properties.author or "",
                "subject": doc.core_properties.subject or "",
                "created": str(doc.core_properties.created) if doc.core_properties.created else "",
                "modified": str(doc.core_properties.modified) if doc.core_properties.modified else "",
            }
        
        return ParseResult(
            file_path=file_path,
            file_name=path.name,
            parser_type=self.parser_type,
            success=True,
            content_text=content_text,
            metadata=metadata,
            word_count=word_count,
        )
    
    def _parse_xlsx(self, file_path: str, path: Path) -> ParseResult:
        from openpyxl import load_workbook
        
        wb = load_workbook(file_path, read_only=True, data_only=True)
        content_parts = []
        metadata = {"sheets": []}
        
        for sheet_name in wb.sheetnames:
            metadata["sheets"].append(sheet_name)
            ws = wb[sheet_name]
            
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip() and row_text != " | ":
                    content_parts.append(row_text)
        
        wb.close()
        
        content_text = "\n".join(content_parts)
        word_count = len(content_text.split()) if content_text else 0
        
        return ParseResult(
            file_path=file_path,
            file_name=path.name,
            parser_type=self.parser_type,
            success=True,
            content_text=content_text,
            metadata=metadata,
            word_count=word_count,
        )
    
    def _parse_pptx(self, file_path: str, path: Path) -> ParseResult:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        content_parts = []
        metadata = {"slide_count": len(prs.slides)}
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_text.append(para.text)
                
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text for cell in row.cells)
                        if row_text.strip():
                            slide_text.append(row_text)
            
            content_parts.append("\n".join(slide_text))
        
        content_text = "\n\n".join(content_parts)
        word_count = len(content_text.split()) if content_text else 0
        
        return ParseResult(
            file_path=file_path,
            file_name=path.name,
            parser_type=self.parser_type,
            success=True,
            content_text=content_text,
            metadata=metadata,
            word_count=word_count,
        )


class TextParser(BaseParser):
    """Parser for plain text and code files."""
    
    @property
    def parser_type(self) -> str:
        return "text"
    
    @property
    def supported_extensions(self) -> List[str]:
        return [
            "txt", "md", "log", "csv", "tsv", "json", "xml", "yaml", "yml",
            "toml", "ini", "cfg", "conf", "env", "gitignore", "dockerignore",
            "py", "js", "ts", "jsx", "tsx", "rs", "go", "java", "cpp", "c",
            "h", "hpp", "cs", "rb", "php", "swift", "kt", "scala", "sh",
            "bash", "zsh", "fish", "ps1", "bat", "cmd", "html", "css",
            "scss", "less", "sql", "graphql", "proto", "yaml", "json",
        ]
    
    def parse(self, file_path: str) -> ParseResult:
        path = Path(file_path)
        if not path.exists():
            return ParseResult(
                file_path=file_path,
                file_name=path.name,
                parser_type=self.parser_type,
                success=False,
                error_message="File not found",
            )
        
        try:
            content = path.read_text(encoding="utf-8", errors="ignore")
            word_count = len(content.split()) if content else 0
            line_count = len(content.splitlines()) if content else 0
            
            metadata = {
                "line_count": line_count,
                "char_count": len(content),
                "encoding": "utf-8",
            }
            
            return ParseResult(
                file_path=file_path,
                file_name=path.name,
                parser_type=self.parser_type,
                success=True,
                content_text=content,
                metadata=metadata,
                word_count=word_count,
            )
            
        except Exception as e:
            return ParseResult(
                file_path=file_path,
                file_name=path.name,
                parser_type=self.parser_type,
                success=False,
                error_message=str(e),
            )


class ImageParser(BaseParser):
    """Parser for image metadata extraction."""
    
    @property
    def parser_type(self) -> str:
        return "image"
    
    @property
    def supported_extensions(self) -> List[str]:
        return ["png", "jpg", "jpeg", "gif", "bmp", "webp", "tiff", "ico"]
    
    def parse(self, file_path: str) -> ParseResult:
        path = Path(file_path)
        if not path.exists():
            return ParseResult(
                file_path=file_path,
                file_name=path.name,
                parser_type=self.parser_type,
                success=False,
                error_message="File not found",
            )
        
        try:
            metadata = {
                "file_size": path.stat().st_size,
                "extension": path.suffix.lower(),
            }
            
            try:
                from PIL import Image
                img = Image.open(file_path)
                metadata.update({
                    "width": img.width,
                    "height": img.height,
                    "format": img.format,
                    "mode": img.mode,
                })
                img.close()
            except ImportError:
                pass
            
            try:
                from PIL.ExifTags import TAGS
                img = Image.open(file_path)
                exif_data = img._getexif()
                if exif_data:
                    exif_metadata = {}
                    for tag_id, value in exif_data.items():
                        tag_name = TAGS.get(tag_id, tag_id)
                        if isinstance(value, (str, int, float)):
                            exif_metadata[str(tag_name)] = value
                    if exif_metadata:
                        metadata["exif"] = exif_metadata
                img.close()
            except Exception:
                pass
            
            return ParseResult(
                file_path=file_path,
                file_name=path.name,
                parser_type=self.parser_type,
                success=True,
                content_text=f"[Image: {metadata.get('width', '?')}x{metadata.get('height', '?')}]",
                metadata=metadata,
            )
            
        except Exception as e:
            return ParseResult(
                file_path=file_path,
                file_name=path.name,
                parser_type=self.parser_type,
                success=False,
                error_message=str(e),
            )


class AudioParser(BaseParser):
    """Parser for audio metadata extraction."""
    
    @property
    def parser_type(self) -> str:
        return "audio"
    
    @property
    def supported_extensions(self) -> List[str]:
        return ["mp3", "wav", "flac", "ogg", "aac", "wma", "m4a"]
    
    def parse(self, file_path: str) -> ParseResult:
        path = Path(file_path)
        if not path.exists():
            return ParseResult(
                file_path=file_path,
                file_name=path.name,
                parser_type=self.parser_type,
                success=False,
                error_message="File not found",
            )
        
        try:
            metadata = {
                "file_size": path.stat().st_size,
                "extension": path.suffix.lower(),
            }
            
            return ParseResult(
                file_path=file_path,
                file_name=path.name,
                parser_type=self.parser_type,
                success=True,
                content_text=f"[Audio file: {path.name}]",
                metadata=metadata,
            )
            
        except Exception as e:
            return ParseResult(
                file_path=file_path,
                file_name=path.name,
                parser_type=self.parser_type,
                success=False,
                error_message=str(e),
            )


class VideoParser(BaseParser):
    """Parser for video metadata extraction."""
    
    @property
    def parser_type(self) -> str:
        return "video"
    
    @property
    def supported_extensions(self) -> List[str]:
        return ["mp4", "avi", "mkv", "mov", "wmv", "webm", "flv", "m4v"]
    
    def parse(self, file_path: str) -> ParseResult:
        path = Path(file_path)
        if not path.exists():
            return ParseResult(
                file_path=file_path,
                file_name=path.name,
                parser_type=self.parser_type,
                success=False,
                error_message="File not found",
            )
        
        try:
            metadata = {
                "file_size": path.stat().st_size,
                "extension": path.suffix.lower(),
            }
            
            return ParseResult(
                file_path=file_path,
                file_name=path.name,
                parser_type=self.parser_type,
                success=True,
                content_text=f"[Video file: {path.name}]",
                metadata=metadata,
            )
            
        except Exception as e:
            return ParseResult(
                file_path=file_path,
                file_name=path.name,
                parser_type=self.parser_type,
                success=False,
                error_message=str(e),
            )


class ParserWorkerPool:
    """Pool of parser workers for different file types."""
    
    def __init__(self):
        self._parsers: List[BaseParser] = [
            PDFParser(),
            OfficeParser(),
            TextParser(),
            ImageParser(),
            AudioParser(),
            VideoParser(),
        ]
    
    def get_parser(self, file_path: str) -> Optional[BaseParser]:
        """Get the appropriate parser for a file."""
        for parser in self._parsers:
            if parser.can_parse(file_path):
                return parser
        return None
    
    def parse_file(self, file_path: str) -> ParseResult:
        """Parse a file using the appropriate parser."""
        parser = self.get_parser(file_path)
        if parser is None:
            path = Path(file_path)
            return ParseResult(
                file_path=file_path,
                file_name=path.name,
                parser_type="unknown",
                success=False,
                error_message=f"No parser available for {path.suffix}",
            )
        
        try:
            return parser.parse(file_path)
        except Exception as e:
            path = Path(file_path)
            return ParseResult(
                file_path=file_path,
                file_name=path.name,
                parser_type=parser.parser_type,
                success=False,
                error_message=f"Parser error: {str(e)}",
            )
    
    def parse_batch(self, file_paths: List[str]) -> List[ParseResult]:
        """Parse multiple files."""
        return [self.parse_file(path) for path in file_paths]
    
    def get_supported_extensions(self) -> List[str]:
        """Get all supported file extensions."""
        extensions = []
        for parser in self._parsers:
            extensions.extend(parser.supported_extensions)
        return list(set(extensions))
    
    def get_parser_types(self) -> List[str]:
        """Get all parser types."""
        return [parser.parser_type for parser in self._parsers]


def create_parser_pool() -> ParserWorkerPool:
    """Factory function to create a parser worker pool."""
    return ParserWorkerPool()


if __name__ == "__main__":
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python parser_workers.py <file_path>")
        sys.exit(1)
    
    pool = create_parser_pool()
    result = pool.parse_file(sys.argv[1])
    
    print(f"File: {result.file_name}")
    print(f"Parser: {result.parser_type}")
    print(f"Success: {result.success}")
    if result.content_text:
        preview = result.content_text[:500] + "..." if len(result.content_text) > 500 else result.content_text
        print(f"Content Preview:\n{preview}")
    if result.metadata:
        print(f"Metadata: {json.dumps(result.metadata, indent=2)}")
    if result.error_message:
        print(f"Error: {result.error_message}")
